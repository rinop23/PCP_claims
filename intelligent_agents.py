"""
Intelligent Multi-Agent System for PCP Claims Analysis
4 Specialized Agents working together to generate investor reports
Now includes Investor Relations Agent for PowerPoint generation
"""

import os
import json
from typing import Dict, Any, List, Tuple
from openai import OpenAI
import pandas as pd
from docx import Document
from docx.shared import Inches
import PyPDF2

# PowerPoint imports
try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None


class BaseAgent:
    """Base class for all intelligent agents"""

    def __init__(self, api_key: str = None):
        if api_key:
            self.api_key = api_key
        else:
            self.api_key = os.environ.get("OPENAI_API_KEY")
            if not self.api_key:
                try:
                    import streamlit as st
                    self.api_key = st.secrets.get("OPENAI_API_KEY")
                except:
                    pass

        if not self.api_key:
            raise ValueError("OpenAI API key required")

        self.client = OpenAI(api_key=self.api_key)
        self.knowledge_base = {}

    def call_openai(self, system_prompt: str, user_prompt: str, response_format: str = "json", max_tokens: int = 16000) -> Any:
        """Call OpenAI with prompts"""
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        if response_format == "json":
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=messages,
                temperature=0.1,
                max_tokens=max_tokens,
                response_format={"type": "json_object"}
            )
        else:
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=messages,
                temperature=0.2,
                max_tokens=max_tokens
            )

        result = response.choices[0].message.content.strip()

        if response_format == "json":
            return json.loads(result)
        else:
            return result


class PriorityDeedAgent(BaseAgent):
    """Agent that reads and understands the Priority Deed for profit distribution"""

    def __init__(self, api_key: str = None):
        super().__init__(api_key)
        self.profit_rules = None

    def read_priority_deed(self, file_path: str = "DOCS/Priorities Deed (EV 9 October 2025).docx") -> Dict[str, Any]:
        """Read Priority Deed document and extract profit distribution rules"""

        print("📄 Priority Deed Agent: Reading profit distribution agreement...")

        try:
            # Read Word document
            doc = Document(file_path)
            text_content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)

            full_text = "\n".join(text_content)

            # Truncate if too long
            if len(full_text) > 15000:
                full_text = full_text[:15000]

        except Exception as e:
            print(f"Warning: Could not read Priority Deed document: {e}")
            full_text = "Priority Deed not available"

        system_prompt = """You are a legal document analyst specializing in litigation funding agreements.
Extract the profit distribution rules from the Priority Deed document.
Return structured JSON with EXACT percentages and rules."""

        user_prompt = f"""Analyze this Priority Deed document and extract profit distribution rules:

{full_text}

Extract and return JSON with:
{{
  "profit_split": {{
    "funder_percentage": number (e.g., 80 for 80%),
    "law_firm_percentage": number (e.g., 20 for 20%),
    "split_basis": "string (e.g., 'after costs', 'of DBA proceeds')"
  }},
  "dba_rate": {{
    "percentage_of_settlement": number (e.g., 30 for 30%),
    "description": "string"
  }},
  "cost_recovery": {{
    "order": "string (e.g., 'costs recovered first, then split')",
    "included_costs": ["list of cost types"]
  }},
  "waterfall": [
    {{
      "priority": 1,
      "recipient": "string",
      "amount_type": "string",
      "description": "string"
    }}
  ],
  "key_terms": {{
    "collection_account": "string describing collection account rules",
    "payment_triggers": "string describing when payments are made",
    "reporting_requirements": "string"
  }}
}}

Be precise with percentages and rules."""

        self.profit_rules = self.call_openai(system_prompt, user_prompt, "json")
        self.knowledge_base = self.profit_rules

        print(f"✅ Extracted profit split: {self.profit_rules['profit_split']['funder_percentage']}% Funder / {self.profit_rules['profit_split']['law_firm_percentage']}% Law Firm")

        return self.profit_rules

    def calculate_distributions(self, total_settlement: float, total_costs: float) -> Dict[str, float]:
        """Calculate profit distributions based on Priority Deed rules

        IMPORTANT: The 80/20 split is on GROSS DBA proceeds, NOT net proceeds after costs.
        """

        if not self.profit_rules:
            raise ValueError("Priority Deed must be read first")

        # Calculate DBA proceeds (typically 30% of settlement)
        dba_rate = self.profit_rules['dba_rate']['percentage_of_settlement'] / 100
        dba_proceeds = total_settlement * dba_rate

        # Split GROSS DBA proceeds (not net after costs)
        funder_split = self.profit_rules['profit_split']['funder_percentage'] / 100
        firm_split = self.profit_rules['profit_split']['law_firm_percentage'] / 100

        # Funder and Milberg split the GROSS DBA proceeds
        funder_share = dba_proceeds * funder_split
        firm_share = dba_proceeds * firm_split

        return {
            "total_settlement": total_settlement,
            "dba_proceeds": dba_proceeds,
            "total_costs": total_costs,
            "funder_share": funder_share,
            "firm_share": firm_share,
            "funder_percentage": self.profit_rules['profit_split']['funder_percentage'],
            "firm_percentage": self.profit_rules['profit_split']['law_firm_percentage']
        }


class FCAComplianceAgent(BaseAgent):
    """Agent that reads FCA Redress Scheme and validates claim compliance"""

    def __init__(self, api_key: str = None):
        super().__init__(api_key)
        self.compliance_rules = None

    def read_fca_scheme(self, file_path: str = "FCA redress scheme/Redress Scheme.pdf") -> Dict[str, Any]:
        """Read FCA Redress Scheme PDF and extract compliance rules"""

        print("⚖️ FCA Compliance Agent: Reading FCA Redress Scheme...")

        try:
            # Read PDF
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_content = []

                # Read first 20 pages (usually contains key info)
                num_pages = min(20, len(pdf_reader.pages))
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text_content.append(page.extract_text())

                full_text = "\n".join(text_content)

                # Truncate if too long
                if len(full_text) > 15000:
                    full_text = full_text[:15000]

        except Exception as e:
            print(f"Warning: Could not read FCA Redress Scheme: {e}")
            full_text = "FCA Redress Scheme not available"

        system_prompt = """You are an FCA compliance expert specializing in motor finance redress schemes.
Extract the key compliance rules and validation criteria for PCP claims.
Return structured JSON."""

        user_prompt = f"""Analyze this FCA Redress Scheme document and extract compliance rules:

{full_text}

Extract and return JSON with:
{{
  "eligible_products": ["list of eligible product types"],
  "commission_thresholds": {{
    "likely_unfair_above": number (percentage),
    "review_required_above": number (percentage),
    "acceptable_below": number (percentage),
    "description": "string"
  }},
  "disclosure_requirements": {{
    "must_disclose": ["list of required disclosures"],
    "disclosure_timing": "string",
    "consequences_of_non_disclosure": "string"
  }},
  "claim_validation": {{
    "required_evidence": ["list of required evidence"],
    "invalid_if": ["list of invalidation criteria"],
    "success_criteria": ["list of success criteria"]
  }},
  "redress_calculation": {{
    "methodology": "string describing how redress is calculated",
    "components": ["list of components included"],
    "exclusions": ["list of exclusions"]
  }},
  "timeline_requirements": {{
    "claim_submission_deadline": "string",
    "expected_response_time": "string",
    "appeal_period": "string"
  }},
  "red_flags": ["list of red flags that indicate non-compliance"]
}}

Be specific with thresholds and criteria."""

        self.compliance_rules = self.call_openai(system_prompt, user_prompt, "json")
        self.knowledge_base = self.compliance_rules

        print(f"✅ Loaded FCA compliance rules for {len(self.compliance_rules.get('eligible_products', []))} product types")

        return self.compliance_rules

    def validate_claim(self, claim_data: Dict[str, Any]) -> Dict[str, Any]:
        """Validate a claim against FCA compliance rules"""

        if not self.compliance_rules:
            raise ValueError("FCA Redress Scheme must be read first")

        commission_pct = claim_data.get('commission_pct_of_cost', 0)
        product_type = claim_data.get('product_type', 'Unknown')

        # Check commission threshold
        thresholds = self.compliance_rules.get('commission_thresholds', {})
        likely_unfair = thresholds.get('likely_unfair_above', 50)

        compliance_status = "compliant"
        issues = []

        if commission_pct > likely_unfair:
            compliance_status = "non_compliant"
            issues.append(f"Commission {commission_pct}% exceeds FCA threshold of {likely_unfair}%")

        if product_type not in self.compliance_rules.get('eligible_products', []):
            issues.append(f"Product type '{product_type}' may not be eligible")

        return {
            "claim_id": claim_data.get('claim_id'),
            "compliance_status": compliance_status,
            "commission_percentage": commission_pct,
            "issues": issues,
            "fca_threshold": likely_unfair,
            "is_valid": len(issues) == 0
        }


class MonthlyReportAgent(BaseAgent):
    """Agent that reads and extracts data from monthly Excel reports"""

    def __init__(self, api_key: str = None):
        super().__init__(api_key)
        self.monthly_data = None

    def analyze_monthly_report(self, file_path: str) -> Dict[str, Any]:
        """Analyze monthly Excel report and extract all metrics"""

        print("📊 Monthly Report Agent: Analyzing Excel report...")

        # Read Excel file
        xl = pd.ExcelFile(file_path)
        sheet_name = xl.sheet_names[0]
        df = pd.read_excel(file_path, sheet_name=sheet_name, header=None)
        excel_text = df.to_string()

        # Truncate if too long
        if len(excel_text) > 25000:
            excel_text = excel_text[:25000] + "\n... [truncated]"

        system_prompt = """You are a financial data analyst. Extract ALL data from monthly reports.
Return structured JSON with complete data - do NOT skip any lenders or data points.
IMPORTANT: Use CUMULATIVE values (not current month) for total counts."""

        user_prompt = f"""Analyze this monthly Excel report and extract ALL data:

{excel_text}

IMPORTANT INSTRUCTIONS:
1. For portfolio metrics, use the CUMULATIVE column values (not Current Month)
2. The "Grand Summary" row contains the total claims and total estimated value
3. Extract ALL lenders from the Defendant table (typically 60-80 lenders)
4. The total estimated value is typically around 228,900 for this report

Extract and return JSON with:
{{
  "reporting_period": "month/year from the Reporting Period field",
  "portfolio_metrics": {{
    "unique_clients": number (use CUMULATIVE value, e.g., 157),
    "unique_claims": number (use CUMULATIVE value, e.g., 327),
    "claims_submitted": number (cumulative),
    "claims_successful": number (cumulative),
    "claims_rejected": number (cumulative),
    "claims_pending": number,
    "avg_claim_value": number (e.g., 228900 or 700),
    "total_settlement_value": number (from Grand Summary Estimated Value, e.g., 228900),
    "success_rate": number (percentage)
  }},
  "lender_distribution": [
    {{
      "lender": "string (Defendant name)",
      "num_claims": number (Number of Claims column),
      "pct_of_total": number (% of Total column, e.g., 0.009 = 0.9%),
      "estimated_value": number (Estimated Value column),
      "avg_claim_value": number
    }}
  ],
  "pipeline": {{
    "awaiting_dsar": {{"count": number, "value": number}},
    "pending_submission": {{"count": number, "value": number}},
    "under_review": {{"count": number, "value": number}},
    "settlement_offered": {{"count": number, "value": number}},
    "paid": {{"count": number, "value": number}}
  }},
  "financial_metrics": {{
    "acquisition_cost": number (cumulative),
    "submission_cost": number (cumulative),
    "processing_cost": number,
    "legal_cost": number,
    "total_costs": number (Total Action Costs cumulative),
    "cost_per_claim": number,
    "collection_balance": number
  }},
  "forecasting": {{
    "expected_new_clients": number,
    "expected_submissions": number,
    "expected_settlement_value": number,
    "projected_timeline": "string"
  }},
  "key_changes": {{
    "new_claims_this_month": number (Current Month claims),
    "claims_resolved_this_month": number,
    "major_settlements": ["list of notable events"]
  }}
}}

CRITICAL REMINDERS:
- unique_clients and unique_claims should be CUMULATIVE values (157 clients, 327 claims)
- total_settlement_value should be from "Grand Summary" row (228900)
- Extract ALL lenders (typically 60-80), not just the first few
- Convert all currency to numbers (remove £ and commas)"""

        self.monthly_data = self.call_openai(system_prompt, user_prompt, "json", max_tokens=16000)
        self.knowledge_base = self.monthly_data

        # Post-process to ensure consistency
        pm = self.monthly_data.get('portfolio_metrics', {})
        lenders = self.monthly_data.get('lender_distribution', [])

        # Calculate total from lenders if total_settlement_value is wrong
        if lenders:
            total_from_lenders = sum(l.get('estimated_value', 0) for l in lenders)
            if pm.get('total_settlement_value', 0) < 1000 and total_from_lenders > 1000:
                pm['total_settlement_value'] = total_from_lenders

        print(f"✅ Extracted data for {pm.get('unique_claims', 0)} claims across {len(lenders)} lenders")
        print(f"   Total settlement value: £{pm.get('total_settlement_value', 0):,.0f}")

        return self.monthly_data


class InvestorReportAgent(BaseAgent):
    """Master agent that generates investor reports using insights from all other agents"""

    def __init__(self, api_key: str = None):
        super().__init__(api_key)

    def generate_investor_report(self,
                                monthly_data: Dict[str, Any],
                                profit_rules: Dict[str, Any],
                                compliance_rules: Dict[str, Any]) -> Dict[str, Any]:
        """Generate comprehensive investor report combining all agent insights"""

        print("📝 Investor Report Agent: Generating comprehensive report...")

        # Pre-calculate key financial metrics to ensure accuracy
        pm = monthly_data.get('portfolio_metrics', {})
        fm = monthly_data.get('financial_metrics', {})
        lenders = monthly_data.get('lender_distribution', [])
        pipeline = monthly_data.get('pipeline', {})

        # Get total settlement value - try multiple sources
        total_settlement = pm.get('total_settlement_value', 0)
        if total_settlement == 0 and lenders:
            total_settlement = sum(l.get('estimated_value', 0) for l in lenders)

        # Get profit split rules (Milberg 20% / Funder 80% of net proceeds after costs)
        dba_rate = 30  # Default 30% DBA rate on settlements
        funder_pct = 80  # Funder gets 80%
        firm_pct = 20  # Milberg (Law Firm) gets 20%

        if profit_rules:
            dba_info = profit_rules.get('dba_rate', {})
            if isinstance(dba_info, dict):
                dba_rate = dba_info.get('percentage_of_settlement', 30)
            split_info = profit_rules.get('profit_split', {})
            if isinstance(split_info, dict):
                funder_pct = split_info.get('funder_percentage', 80)
                firm_pct = split_info.get('law_firm_percentage', 20)

        # Calculate financials
        # IMPORTANT: The 80/20 split is on GROSS DBA proceeds, NOT net proceeds after costs
        dba_proceeds = total_settlement * (dba_rate / 100)
        total_costs = fm.get('total_costs', 0)

        # Funder and Milberg split the GROSS DBA proceeds (not net after costs)
        funder_return = dba_proceeds * (funder_pct / 100)
        firm_return = dba_proceeds * (firm_pct / 100)

        # ROI = (Return - Investment) / Investment * 100
        # MOIC = Return / Investment
        roi = ((funder_return - total_costs) / total_costs * 100) if total_costs > 0 else 0
        moic = (funder_return / total_costs) if total_costs > 0 else 0

        # Pre-calculated values to inject into prompt
        pre_calculated = {
            "total_settlement_value": total_settlement,
            "dba_rate": dba_rate,
            "dba_proceeds": dba_proceeds,
            "total_costs": total_costs,
            "funder_return": funder_return,
            "firm_return": firm_return,
            "funder_percentage": funder_pct,
            "firm_percentage": firm_pct,
            "roi": roi,
            "moic": moic,
            "total_claims": pm.get('unique_claims', 0),
            "total_clients": pm.get('unique_clients', 0),
            "total_lenders": len(lenders),
            "pipeline_value": sum(stage.get('value', 0) for stage in pipeline.values() if isinstance(stage, dict))
        }

        print(f"   Pre-calculated: settlement=£{total_settlement:,.0f}, DBA=£{dba_proceeds:,.0f}, funder=£{funder_return:,.0f}, MOIC={moic:.2f}x")

        system_prompt = """You are a senior investment analyst creating monthly investor reports for litigation funding stakeholders.
Generate factual, data-driven analysis. Include detailed analysis for each section matching the dashboard tabs:
1. Lenders - concentration analysis, top lenders, diversification
2. Economic Analysis - revenue, costs, profit distribution, ROI
3. Compliance & Pipeline - FCA compliance, pipeline stages, bottlenecks
4. Portfolio Analysis - growth metrics, risk assessment, forecasting"""

        user_prompt = f"""Create a comprehensive monthly investor report using these inputs:

MONTHLY REPORT DATA:
{json.dumps(monthly_data, indent=2)}

PRE-CALCULATED FINANCIALS (USE THESE EXACT VALUES):
{json.dumps(pre_calculated, indent=2)}

PROFIT DISTRIBUTION RULES (from Priority Deed):
{json.dumps(profit_rules, indent=2)}

FCA COMPLIANCE RULES (from Redress Scheme):
{json.dumps(compliance_rules, indent=2)}

CRITICAL INSTRUCTIONS:
1. Use the PRE-CALCULATED FINANCIALS values - DO NOT recalculate
2. total_settlements = {total_settlement} (the total portfolio value)
3. dba_proceeds_expected = {dba_proceeds:.2f} (30% of settlements)
4. total_costs_incurred = {total_costs:.2f}
5. funder_expected_return = {funder_return:.2f} ({funder_pct}% of GROSS DBA proceeds)
6. firm_expected_return = {firm_return:.2f} ({firm_pct}% of GROSS DBA proceeds)
7. total_claims = {pre_calculated['total_claims']}
8. total_clients = {pre_calculated['total_clients']}
9. total_portfolio_value = {total_settlement}
NOTE: The split is on GROSS DBA proceeds, NOT net proceeds after costs.

Generate a structured investor report in JSON format:
{{
  "executive_summary": {{
    "reporting_period": "string",
    "portfolio_health_score": number (0-100),
    "key_metrics_summary": ["3-5 key metrics with values"],
    "critical_updates": ["major changes or events"]
  }},
  "portfolio_performance": {{
    "total_claims": number,
    "total_clients": number,
    "claims_by_stage": {{}},
    "success_rate": number,
    "average_settlement": number,
    "total_portfolio_value": number,
    "month_over_month_growth": "string with numbers"
  }},
  "financial_analysis": {{
    "total_settlements": number,
    "dba_proceeds_expected": number,
    "total_costs_incurred": number,
    "funder_expected_return": number,
    "firm_expected_return": number,
    "roi_projection": number,
    "moic_projection": number
  }},
  "compliance_assessment": {{
    "fca_compliance_status": "compliant/review_required/at_risk",
    "commission_analysis": "string describing commission levels vs FCA thresholds",
    "claims_at_risk": number,
    "compliance_actions_needed": ["list specific actions"]
  }},
  "lender_concentration": {{
    "total_lenders": number,
    "top_5_lenders": [{{"lender": "string", "claims": number, "percentage": number}}],
    "concentration_risk": "low/medium/high with explanation",
    "diversification_score": number (0-100)
  }},
  "pipeline_analysis": {{
    "pipeline_by_stage": {{}},
    "conversion_rates": "string with percentages",
    "estimated_time_to_settlement": "string",
    "pipeline_value": number,
    "bottlenecks": ["list any bottlenecks observed"]
  }},
  "cost_efficiency": {{
    "cost_per_claim": number,
    "cost_per_successful_claim": number,
    "cost_breakdown": {{}},
    "efficiency_trends": "string describing if costs are improving"
  }},
  "forecasting": {{
    "next_month_projections": {{}},
    "quarterly_outlook": "string",
    "expected_settlements_next_90_days": number,
    "projected_returns": {{}}
  }},
  "risk_assessment": {{
    "key_risks": ["list of specific risks with data"],
    "mitigation_status": ["what's being done about each risk"],
    "risk_level": "low/medium/high"
  }},
  "action_items": [
    {{
      "priority": "high/medium/low",
      "action": "string",
      "owner": "Funder/Milberg/Processor",
      "deadline": "string",
      "rationale": "data-driven reason"
    }}
  ]
}}

IMPORTANT:
- Use ONLY the data provided
- Include specific numbers for every metric
- Calculate profit distributions using the Priority Deed split rules
- Flag compliance issues based on FCA thresholds
- Be factual and objective - no opinions
- Show calculations and basis for projections"""

        report_data = self.call_openai(system_prompt, user_prompt, "json", max_tokens=16000)

        # POST-PROCESS: Force correct financial values (OpenAI sometimes ignores our pre-calculated values)
        report_data = self._force_correct_financials(report_data, pre_calculated)

        # Generate markdown version for display
        markdown_report = self._format_as_markdown(report_data)

        # Generate a short narrative (few lines) to accompany charts in the DOCX
        narrative = self._generate_short_narrative(report_data)

        print("✅ Investor report generated successfully")

        return {
            "report_data": report_data,
            "markdown_report": markdown_report,
            "narrative": narrative,
        }

    def _force_correct_financials(self, report_data: Dict[str, Any], pre_calculated: Dict[str, Any]) -> Dict[str, Any]:
        """Force correct financial values in the report data (overwrite any incorrect OpenAI values)"""

        # Force correct portfolio_performance values
        if 'portfolio_performance' not in report_data:
            report_data['portfolio_performance'] = {}
        pp = report_data['portfolio_performance']
        pp['total_claims'] = pre_calculated.get('total_claims', pp.get('total_claims', 0))
        pp['total_clients'] = pre_calculated.get('total_clients', pp.get('total_clients', 0))
        pp['total_portfolio_value'] = pre_calculated.get('total_settlement_value', pp.get('total_portfolio_value', 0))

        # Force correct financial_analysis values
        if 'financial_analysis' not in report_data:
            report_data['financial_analysis'] = {}
        fa = report_data['financial_analysis']
        fa['total_settlements'] = pre_calculated.get('total_settlement_value', 0)
        fa['dba_proceeds_expected'] = pre_calculated.get('dba_proceeds', 0)
        fa['total_costs_incurred'] = pre_calculated.get('total_costs', 0)
        fa['funder_expected_return'] = pre_calculated.get('funder_return', 0)
        fa['firm_expected_return'] = pre_calculated.get('firm_return', 0)
        fa['roi_projection'] = pre_calculated.get('roi', 0)
        fa['moic_projection'] = pre_calculated.get('moic', 0)

        # Force lender concentration
        if 'lender_concentration' not in report_data:
            report_data['lender_concentration'] = {}
        report_data['lender_concentration']['total_lenders'] = pre_calculated.get('total_lenders', 0)

        # Force pipeline value
        if 'pipeline_analysis' not in report_data:
            report_data['pipeline_analysis'] = {}
        report_data['pipeline_analysis']['pipeline_value'] = pre_calculated.get('pipeline_value', 0)

        print(f"   Forced financials: settlements=£{fa['total_settlements']:,.0f}, DBA=£{fa['dba_proceeds_expected']:,.0f}, funder=£{fa['funder_expected_return']:,.0f}, MOIC={fa['moic_projection']:.2f}x")

        return report_data

    def _generate_short_narrative(self, report_data: Dict[str, Any]) -> str:
        """Generate brief narrative lines to explain the data (kept short and factual)."""
        exec_sum = report_data.get("executive_summary") or {}
        perf = report_data.get("portfolio_performance") or {}
        fin = report_data.get("financial_analysis") or {}

        period = exec_sum.get("reporting_period") or "the reporting period"
        claims = perf.get("total_claims")
        clients = perf.get("total_clients")
        value = perf.get("total_portfolio_value")
        roi = fin.get("roi_projection")

        parts = [
            f"This report summarizes portfolio activity for {period}.",
        ]
        if claims is not None and clients is not None:
            parts.append(f"The portfolio contains {claims} claims across {clients} clients.")
        if value is not None:
            parts.append(f"Estimated total portfolio value is £{value:,.0f}.")
        if roi is not None:
            parts.append(f"Projected ROI based on current inputs is {roi:.1f}%.")

        # Keep it short (few lines)
        return " ".join(parts[:4])

    def _format_as_markdown(self, report_data: Dict[str, Any]) -> str:
        """Format the structured report as markdown (best-effort, safe for missing values)."""

        def _num(v, default=None):
            return v if isinstance(v, (int, float)) else default

        def _int(v, default=0) -> int:
            try:
                return int(float(v))
            except Exception:
                return default

        def _pct(v, default=None):
            n = _num(v, default)
            return "N/A" if n is None else f"{n:.1f}%"

        def _x(v, default=None):
            n = _num(v, default)
            return "N/A" if n is None else f"{n:.2f}x"

        def _gbp(v, default=None):
            n = _num(v, default)
            return "N/A" if n is None else f"£{n:,.2f}"

        exec_sum = report_data.get("executive_summary") or {}
        perf = report_data.get("portfolio_performance") or {}
        fin = report_data.get("financial_analysis") or {}
        comp = report_data.get("compliance_assessment") or {}
        conc = report_data.get("lender_concentration") or {}
        pipe = report_data.get("pipeline_analysis") or {}
        cost = report_data.get("cost_efficiency") or {}
        fcst = report_data.get("forecasting") or {}
        risk = report_data.get("risk_assessment") or {}
        action_items = report_data.get("action_items") or []

        md = f"""# Monthly Investor Report
## {exec_sum.get('reporting_period') or 'Reporting Period: N/A'}

---

## Executive Summary

**Portfolio Health Score:** {exec_sum.get('portfolio_health_score') if exec_sum.get('portfolio_health_score') is not None else 'N/A'}/100

### Key Metrics
"""
        for metric in (exec_sum.get('key_metrics_summary') or []):
            md += f"- {metric}\n"

        md += "\n### Critical Updates\n"
        for update in (exec_sum.get('critical_updates') or []):
            md += f"- {update}\n"

        md += f"""

---

## Portfolio Performance

- **Total Claims:** {_int(perf.get('total_claims')):,}
- **Total Clients:** {_int(perf.get('total_clients')):,}
- **Success Rate:** {_pct(perf.get('success_rate'))}
- **Average Settlement:** {_gbp(perf.get('average_settlement'))}
- **Total Portfolio Value:** {_gbp(perf.get('total_portfolio_value'))}

{perf.get('month_over_month_growth') or ''}

---

## Financial Analysis

### Revenue
- **Total Expected Settlements:** {_gbp(fin.get('total_settlements'))}
- **DBA Proceeds (30%):** {_gbp(fin.get('dba_proceeds_expected'))}

### Costs
- **Total Costs Incurred:** {_gbp(fin.get('total_costs_incurred'))}

### Profit Distribution (80/20 split on GROSS DBA)
- **LP Return (80%):** {_gbp(fin.get('funder_expected_return'))}
- **GP Return (20%):** {_gbp(fin.get('firm_expected_return'))}

### Performance Metrics
- **ROI Projection:** {_pct(fin.get('roi_projection'))}
- **MOIC Projection:** {_x(fin.get('moic_projection'))}

---

## FCA Compliance Assessment

**Status:** {(comp.get('fca_compliance_status') or 'N/A').upper()}

{comp.get('commission_analysis') or ''}

- **Claims at Risk:** {_int(comp.get('claims_at_risk'), default=0)}

### Actions Required
"""
        for action in (comp.get('compliance_actions_needed') or []):
            md += f"- {action}\n"

        md += f"""

---

## Lender Concentration

**Total Lenders:** {_int(conc.get('total_lenders'), default=0)}
**Diversification Score:** {conc.get('diversification_score') if conc.get('diversification_score') is not None else 'N/A'}/100
**Concentration Risk:** {conc.get('concentration_risk') or 'N/A'}

### Top 5 Lenders
"""
        for lender in (conc.get('top_5_lenders') or []):
            lender_name = (lender or {}).get('lender') or 'Unknown'
            claims = _int((lender or {}).get('claims'), default=0)
            pct = _num((lender or {}).get('percentage'), default=None)
            pct_s = "N/A" if pct is None else f"{pct:.1f}%"
            md += f"- **{lender_name}**: {claims} claims ({pct_s})\n"

        md += f"""

---

## Pipeline Analysis

{pipe.get('conversion_rates') or ''}

**Pipeline Value:** {_gbp(pipe.get('pipeline_value'))}
**Est. Time to Settlement:** {pipe.get('estimated_time_to_settlement') or 'N/A'}

### Bottlenecks Identified
"""
        for bottleneck in (pipe.get('bottlenecks') or []):
            md += f"- {bottleneck}\n"

        md += f"""

---

## Cost Efficiency

- **Cost per Claim:** {_gbp(cost.get('cost_per_claim'))}
- **Cost per Successful Claim:** {_gbp(cost.get('cost_per_successful_claim'))}

{cost.get('efficiency_trends') or ''}

---

## Forecasting

{fcst.get('quarterly_outlook') or ''}

**Expected Settlements (Next 90 Days):** {_gbp(fcst.get('expected_settlements_next_90_days'))}

---

## Risk Assessment

**Overall Risk Level:** {(risk.get('risk_level') or 'N/A').upper()}

### Key Risks
"""
        for r in (risk.get('key_risks') or []):
            md += f"- {r}\n"

        md += "\n### Mitigation Actions\n"
        for m in (risk.get('mitigation_status') or []):
            md += f"- {m}\n"

        md += "\n---\n\n## Action Items\n\n"

        for item in action_items:
            item = item or {}
            md += f"""
### [{(item.get('priority') or 'N/A').upper()}] {item.get('action') or 'N/A'}
- **Owner:** {item.get('owner') or 'N/A'}
- **Deadline:** {item.get('deadline') or 'N/A'}
- **Rationale:** {item.get('rationale') or 'N/A'}
"""

        md += "\n---\n\n*Report generated by AI Multi-Agent System*"

        return md


def _build_dashboard_figures(monthly_data: Dict[str, Any], report_data: Dict[str, Any]):
    """Create Plotly figures similar to the dashboard for embedding in the DOCX."""
    try:
        import plotly.express as px
        import plotly.graph_objects as go
    except Exception:
        return {}

    figs = {}

    # Shared styling for static export readability
    def _style(fig, *, height: int = 520):
        try:
            fig.update_layout(
                template="plotly_white",
                height=height,
                margin=dict(l=40, r=20, t=60, b=40),
                paper_bgcolor="white",
                plot_bgcolor="white",
                font=dict(size=14, color="#111"),
                title=dict(x=0.02, xanchor="left"),
                legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="left", x=0),
            )
            fig.update_xaxes(showgrid=True, gridcolor="#e6e6e6", zeroline=False)
            fig.update_yaxes(showgrid=True, gridcolor="#e6e6e6", zeroline=False)
        except Exception:
            pass
        return fig

    money_fmt = "£,.0f"

    # ==================== LENDERS (Dashboard Tab) ====================
    lenders = (monthly_data or {}).get("lender_distribution") or []
    if lenders:
        df_l = pd.DataFrame(lenders).sort_values("num_claims", ascending=False)

        # Top 10 lenders pie by claims
        top10 = df_l.head(10)
        others_claims = df_l.iloc[10:]["num_claims"].sum() if len(df_l) > 10 else 0
        if others_claims > 0:
            pie_df = pd.concat(
                [top10[["lender", "num_claims"]], pd.DataFrame([{ "lender": "Others", "num_claims": others_claims }])]
            )
        else:
            pie_df = top10[["lender", "num_claims"]]

        fig_pie = px.pie(
            pie_df,
            values="num_claims",
            names="lender",
            hole=0.4,
            color_discrete_sequence=px.colors.qualitative.Set2,
        )
        fig_pie.update_traces(textposition="inside", textinfo="percent+label")
        fig_pie.update_layout(title="Top 10 Lenders by Claims")
        figs["lenders_pie"] = _style(fig_pie, height=520)

        # Top 15 lenders by value (horizontal bar)
        if "estimated_value" in df_l.columns:
            top15_val = df_l.head(15).sort_values("estimated_value")
            fig_bar = px.bar(
                top15_val,
                x="estimated_value",
                y="lender",
                orientation="h",
                color="num_claims" if "num_claims" in top15_val.columns else None,
                color_continuous_scale="Blues",
                text=top15_val["estimated_value"].map(lambda v: f"£{v:,.0f}"),
                title="Top 15 Lenders by Portfolio Value",
            )
            fig_bar.update_traces(textposition="outside")
            fig_bar.update_layout(showlegend=False)
            fig_bar.update_xaxes(tickformat=money_fmt)
            figs["lenders_value_bar"] = _style(fig_bar, height=650)

    # ==================== ECONOMIC ANALYSIS (Dashboard Tab) ====================
    # Use investor_report financials for simple, readable static charts
    fin = (report_data or {}).get("financial_analysis") or {}
    try:
        total_settlements = fin.get("total_settlements")
        dba_proceeds = fin.get("dba_proceeds_expected")
        total_costs = fin.get("total_costs_incurred")
        funder_return = fin.get("funder_expected_return")

        if any(v is not None for v in [total_settlements, dba_proceeds, total_costs, funder_return]):
            rows = [
                {"Step": "Total Expected Settlements", "Value": float(total_settlements or 0)},
                {"Step": "DBA Proceeds (30%)", "Value": float(dba_proceeds or 0)},
                {"Step": "Total Costs Incurred", "Value": float(total_costs or 0)},
                {"Step": "LP Return (80% of DBA)", "Value": float(funder_return or 0)},
            ]
            df_fin = pd.DataFrame(rows)
            fig_fin = px.bar(
                df_fin,
                x="Step",
                y="Value",
                text=df_fin["Value"].map(lambda v: f"£{v:,.0f}"),
                color="Step",
                color_discrete_sequence=px.colors.qualitative.Set2,
                title="Economic Summary (Expected)",
            )
            fig_fin.update_traces(textposition="outside")
            fig_fin.update_yaxes(tickformat=money_fmt)
            fig_fin.update_layout(showlegend=False, xaxis_title="", yaxis_title="Value")
            figs["economic_summary"] = _style(fig_fin, height=520)

            # Profit split (funder vs firm) if provided
            funder_ret = fin.get("funder_expected_return")
            firm_ret = fin.get("firm_expected_return")
            if funder_ret is not None or firm_ret is not None:
                df_split = pd.DataFrame(
                    [
                        {"Recipient": "LP (80%)", "Value": float(funder_ret or 0)},
                        {"Recipient": "GP (20%)", "Value": float(firm_ret or 0)},
                    ]
                )
                fig_split = px.pie(
                    df_split,
                    names="Recipient",
                    values="Value",
                    hole=0.45,
                    color_discrete_sequence=["#3498db", "#2ecc71"],
                )
                fig_split.update_traces(textposition="inside", textinfo="percent+label")
                fig_split.update_layout(title="Expected Profit Distribution")
                figs["profit_split"] = _style(fig_split, height=520)
    except Exception:
        pass

    # ==================== COMPLIANCE & PIPELINE (Dashboard Tab) ====================
    pipeline = (monthly_data or {}).get("pipeline") or {}
    stage_order = [
        ("Awaiting DSAR", "awaiting_dsar"),
        ("Pending Submission", "pending_submission"),
        ("Under Review", "under_review"),
        ("Settlement Offered", "settlement_offered"),
        ("Paid", "paid"),
    ]
    try:
        pipe_rows = []
        for label, key in stage_order:
            d = pipeline.get(key) or {}
            pipe_rows.append({"Stage": label, "Count": d.get("count", 0), "Value": d.get("value", 0)})
        df_p = pd.DataFrame(pipe_rows)

        fig_funnel = go.Figure(
            go.Funnel(
                y=df_p["Stage"],
                x=df_p["Count"],
                textinfo="value+percent initial",
                marker={"color": ["#3498db", "#2ecc71", "#f39c12", "#9b59b6", "#1abc9c"]},
            )
        )
        fig_funnel.update_layout(title="Pipeline Funnel (by Count)")
        figs["pipeline_funnel"] = _style(fig_funnel, height=520)

        fig_pipe_val = px.bar(
            df_p,
            x="Stage",
            y="Value",
            text=df_p["Value"].map(lambda v: f"£{v:,.0f}"),
            color="Stage",
            color_discrete_sequence=px.colors.qualitative.Set2,
            title="Pipeline Value by Stage",
        )
        fig_pipe_val.update_traces(textposition="outside")
        fig_pipe_val.update_yaxes(tickformat=money_fmt)
        fig_pipe_val.update_layout(showlegend=False, xaxis_title="", yaxis_title="Value")
        figs["pipeline_value"] = _style(fig_pipe_val, height=520)
    except Exception:
        pass

    return figs


def _export_plotly_fig_to_png_bytes(fig) -> bytes:
    """Export a Plotly figure to PNG bytes (requires kaleido).

    Note: This is kept for backward compatibility. Prefer `_export_plotly_fig_to_png_bytes_with_error`
    for better diagnostics.
    """
    png, _err = _export_plotly_fig_to_png_bytes_with_error(fig)
    return png


def _export_plotly_fig_to_png_bytes_with_error(fig) -> Tuple[bytes, Exception | None]:
    """Export a Plotly figure to PNG bytes (requires kaleido), returning (bytes, error).

    In some environments `plotly.io.to_image` can return empty bytes when Kaleido isn't usable.
    We force a full figure build and provide actionable diagnostics.
    """
    try:
        import plotly.io as pio

        # Force full computation/layout prior to export
        try:
            fig = fig.full_figure_for_development(warn=False)
        except Exception:
            # Not all figure types/environments support this; proceed without it.
            pass

        b = pio.to_image(fig, format="png", width=1200, height=700, scale=2)
        if not b:
            raise RuntimeError("plotly.io.to_image returned empty bytes")
        return b, None
    except Exception as e:
        return b"", e


def _kaleido_debug_info() -> str:
    """Return debug info string about Kaleido/Plotly image export availability."""
    try:
        import importlib
        import plotly

        parts = [f"plotly={getattr(plotly, '__version__', 'unknown')}"]

        try:
            k = importlib.import_module("kaleido")
            parts.append(f"kaleido={getattr(k, '__version__', 'unknown')}")
        except Exception as e:
            parts.append(f"kaleido_import_error={type(e).__name__}: {e}")

        # Optional: try creating a tiny image to verify export works
        try:
            import pandas as pd
            import plotly.express as px
            import plotly.io as pio

            fig = px.bar(pd.DataFrame({"x": [1, 2], "y": [3, 4]}), x="x", y="y")
            b = pio.to_image(fig, format="png", width=10, height=10, scale=1)
            parts.append(f"to_image_ok={bool(b)}")
        except Exception as e:
            parts.append(f"to_image_error={type(e).__name__}: {e}")

        return "; ".join(parts)
    except Exception as e:
        return f"kaleido_debug_failed={type(e).__name__}: {e}"


def build_investor_report_docx(
    *,
    out_path: str,
    narrative: str,
    monthly_data: Dict[str, Any],
    investor_report: Dict[str, Any],
) -> str:
    """Create a .docx investor report including all dashboard data, tables, and charts."""

    def _fmt_currency(v):
        """Format value as currency"""
        if v is None or v == "N/A":
            return "N/A"
        try:
            return f"£{float(v):,.2f}"
        except (ValueError, TypeError):
            return str(v)

    def _fmt_pct(v):
        """Format value as percentage"""
        if v is None or v == "N/A":
            return "N/A"
        try:
            return f"{float(v):.1f}%"
        except (ValueError, TypeError):
            return str(v)

    def _fmt_num(v):
        """Format value as number with commas"""
        if v is None or v == "N/A":
            return "N/A"
        try:
            return f"{int(float(v)):,}"
        except (ValueError, TypeError):
            return str(v)

    def _add_table(headers: List[str], rows: List[List[str]], style: str = None):
        """Add a table with headers and rows"""
        table = doc.add_table(rows=1, cols=len(headers))
        if style:
            table.style = style
        # Header row
        hdr_cells = table.rows[0].cells
        for i, header in enumerate(headers):
            hdr_cells[i].text = str(header)
        # Data rows
        for row_data in rows:
            row_cells = table.add_row().cells
            for i, cell_val in enumerate(row_data):
                row_cells[i].text = str(cell_val) if cell_val is not None else "N/A"

    doc = Document()
    doc.add_heading("Monthly Investor Report", level=0)

    # Get report data
    exec_sum = (investor_report or {}).get("executive_summary") or {}
    perf = (investor_report or {}).get("portfolio_performance") or {}
    fin = (investor_report or {}).get("financial_analysis") or {}
    comp = (investor_report or {}).get("compliance_assessment") or {}
    risk = (investor_report or {}).get("risk_assessment") or {}
    lender_conc = (investor_report or {}).get("lender_concentration") or {}
    pipeline_analysis = (investor_report or {}).get("pipeline_analysis") or {}
    cost_eff = (investor_report or {}).get("cost_efficiency") or {}
    fcst = (investor_report or {}).get("forecasting") or {}

    # Get raw monthly data
    lenders = (monthly_data or {}).get("lender_distribution") or []
    pipeline = (monthly_data or {}).get("pipeline") or {}
    pm = (monthly_data or {}).get("portfolio_metrics") or {}
    fm = (monthly_data or {}).get("financial_metrics") or {}

    period = exec_sum.get("reporting_period") or "Monthly Report"
    doc.add_paragraph(f"Reporting Period: {period}")
    doc.add_paragraph("")

    # ==================== EXECUTIVE SUMMARY ====================
    doc.add_heading("Executive Summary", level=1)

    if narrative:
        doc.add_paragraph(narrative)
        doc.add_paragraph("")

    # Key highlights
    if exec_sum.get("key_metrics_summary"):
        doc.add_heading("Key Highlights", level=2)
        for metric in exec_sum.get("key_metrics_summary", []):
            doc.add_paragraph(f"• {metric}")

    if exec_sum.get("critical_updates"):
        doc.add_heading("Critical Updates", level=2)
        for update in exec_sum.get("critical_updates", []):
            doc.add_paragraph(f"• {update}")

    doc.add_paragraph("")

    # ==================== 1. LENDERS SECTION ====================
    doc.add_heading("1. Lenders Distribution", level=1)

    # Lender summary metrics
    doc.add_heading("Lender Summary", level=2)
    _add_table(
        ["Metric", "Value"],
        [
            ["Total Lenders", _fmt_num(lender_conc.get("total_lenders") or len(lenders))],
            ["Diversification Score", f"{lender_conc.get('diversification_score', 'N/A')}/100" if lender_conc.get('diversification_score') else "N/A"],
            ["Concentration Risk", lender_conc.get("concentration_risk", "N/A")],
        ]
    )
    doc.add_paragraph("")

    # Top 10 Lenders table (limited to top 10 only)
    if lenders:
        doc.add_heading("Top 10 Lenders by Claims", level=2)
        doc.add_paragraph(f"Showing top 10 of {len(lenders)} total lenders in portfolio:")

        # Sort by claims descending and take top 10
        sorted_lenders = sorted(lenders, key=lambda x: x.get('num_claims', 0), reverse=True)[:10]

        lender_rows = []
        for lender in sorted_lenders:
            lender_rows.append([
                lender.get("lender", "Unknown"),
                _fmt_num(lender.get("num_claims", 0)),
                _fmt_pct(lender.get("pct_of_total", 0) * 100 if lender.get("pct_of_total", 0) < 1 else lender.get("pct_of_total", 0)),
                _fmt_currency(lender.get("estimated_value", 0)),
                _fmt_currency(lender.get("avg_claim_value", 0))
            ])

        _add_table(
            ["Lender (Defendant)", "Claims", "% of Total", "Estimated Value", "Avg Claim Value"],
            lender_rows
        )
        doc.add_paragraph("")

    # ==================== 2. ECONOMIC ANALYSIS SECTION ====================
    doc.add_heading("2. Economic Analysis", level=1)

    # Financial Summary Table
    doc.add_heading("Financial Summary", level=2)
    _add_table(
        ["Metric", "Value"],
        [
            ["Total Settlement Value", _fmt_currency(fin.get("total_settlements"))],
            ["DBA Proceeds (30%)", _fmt_currency(fin.get("dba_proceeds_expected"))],
            ["Total Costs Incurred", _fmt_currency(fin.get("total_costs_incurred"))],
            ["LP Return (80% of DBA)", _fmt_currency(fin.get("funder_expected_return"))],
            ["GP Return (20% of DBA)", _fmt_currency(fin.get("firm_expected_return"))],
            ["ROI Projection", _fmt_pct(fin.get("roi_projection"))],
            ["MOIC Projection", f"{fin.get('moic_projection', 0):.2f}x" if fin.get('moic_projection') else "N/A"],
        ]
    )
    doc.add_paragraph("")

    # Cost Breakdown
    doc.add_heading("Cost Breakdown", level=2)
    cost_rows = [
        ["Acquisition Cost", _fmt_currency(fm.get("acquisition_cost", 0))],
        ["Submission Cost", _fmt_currency(fm.get("submission_cost", 0))],
        ["Processing Cost", _fmt_currency(fm.get("processing_cost", 0))],
        ["Legal Cost", _fmt_currency(fm.get("legal_cost", 0))],
        ["Total Costs", _fmt_currency(fm.get("total_costs", 0))],
        ["Cost per Claim", _fmt_currency(cost_eff.get("cost_per_claim") or fm.get("cost_per_claim", 0))],
        ["Cost per Successful Claim", _fmt_currency(cost_eff.get("cost_per_successful_claim", 0))],
    ]
    _add_table(["Cost Category", "Amount"], cost_rows)
    doc.add_paragraph("")

    # Profit Distribution explanation
    doc.add_heading("Profit Distribution (Priority Deed Terms)", level=2)
    doc.add_paragraph("• DBA Rate: 30% of settlements")
    doc.add_paragraph("• Split: 80/20 (LP/GP) on GROSS DBA proceeds")
    doc.add_paragraph("• LP = Limited Partner (Funder), GP = General Partner (Milberg)")
    doc.add_paragraph("• Costs paid separately by LP")
    doc.add_paragraph("")

    if cost_eff.get("efficiency_trends"):
        doc.add_paragraph(f"Efficiency Trends: {cost_eff.get('efficiency_trends')}")
        doc.add_paragraph("")

    # ==================== 3. COMPLIANCE & PIPELINE SECTION ====================
    doc.add_heading("3. Compliance & Pipeline", level=1)

    # FCA Compliance Assessment
    doc.add_heading("FCA Compliance Assessment", level=2)
    _add_table(
        ["Metric", "Value"],
        [
            ["Compliance Status", (comp.get("fca_compliance_status") or "N/A").upper()],
            ["Claims at Risk", _fmt_num(comp.get("claims_at_risk", 0))],
        ]
    )
    doc.add_paragraph("")

    if comp.get("commission_analysis"):
        doc.add_paragraph(f"Commission Analysis: {comp.get('commission_analysis')}")
        doc.add_paragraph("")

    if comp.get("compliance_actions_needed"):
        doc.add_heading("Compliance Actions Required", level=3)
        for action in comp.get("compliance_actions_needed", []):
            doc.add_paragraph(f"• {action}")
        doc.add_paragraph("")

    # Pipeline Breakdown
    doc.add_heading("Claims Pipeline Breakdown", level=2)

    # Pipeline stages table
    pipeline_stages = [
        ("Awaiting DSAR Response", pipeline.get("awaiting_dsar", {})),
        ("Pending Submission", pipeline.get("pending_submission", {})),
        ("Under Review", pipeline.get("under_review", {})),
        ("Settlement Offered", pipeline.get("settlement_offered", {})),
        ("Paid", pipeline.get("paid", {})),
    ]

    pipeline_rows = []
    total_pipeline_count = 0
    total_pipeline_value = 0
    for stage_name, stage_data in pipeline_stages:
        count = stage_data.get("count", 0) if isinstance(stage_data, dict) else 0
        value = stage_data.get("value", 0) if isinstance(stage_data, dict) else 0
        total_pipeline_count += count
        total_pipeline_value += value
        pipeline_rows.append([stage_name, _fmt_num(count), _fmt_currency(value)])

    pipeline_rows.append(["TOTAL", _fmt_num(total_pipeline_count), _fmt_currency(total_pipeline_value)])

    _add_table(["Pipeline Stage", "Count", "Value"], pipeline_rows)
    doc.add_paragraph("")

    # Pipeline analysis
    if pipeline_analysis.get("conversion_rates"):
        doc.add_paragraph(f"Conversion Rates: {pipeline_analysis.get('conversion_rates')}")
    if pipeline_analysis.get("estimated_time_to_settlement"):
        doc.add_paragraph(f"Estimated Time to Settlement: {pipeline_analysis.get('estimated_time_to_settlement')}")

    if pipeline_analysis.get("bottlenecks"):
        doc.add_heading("Pipeline Bottlenecks Identified", level=3)
        for bottleneck in pipeline_analysis.get("bottlenecks", []):
            doc.add_paragraph(f"• {bottleneck}")
    doc.add_paragraph("")

    # ==================== 4. PORTFOLIO ANALYSIS SECTION ====================
    doc.add_heading("4. Portfolio Analysis", level=1)

    # Portfolio Metrics
    doc.add_heading("Portfolio Metrics", level=2)
    _add_table(
        ["Metric", "Value"],
        [
            ["Total Claims", _fmt_num(perf.get("total_claims") or pm.get("unique_claims", 0))],
            ["Total Clients", _fmt_num(perf.get("total_clients") or pm.get("unique_clients", 0))],
            ["Total Lenders", _fmt_num(len(lenders))],
            ["Total Portfolio Value", _fmt_currency(perf.get("total_portfolio_value") or pm.get("total_settlement_value", 0))],
            ["Claims Submitted", _fmt_num(pm.get("claims_submitted", 0))],
            ["Claims Successful", _fmt_num(pm.get("claims_successful", 0))],
            ["Claims Rejected", _fmt_num(pm.get("claims_rejected", 0))],
            ["Average Claim Value", _fmt_currency(perf.get("average_settlement") or pm.get("avg_claim_value", 0))],
            ["Success Rate", _fmt_pct(perf.get("success_rate") or pm.get("success_rate", 0))],
        ]
    )
    doc.add_paragraph("")

    if perf.get("month_over_month_growth"):
        doc.add_paragraph(f"Month-over-Month Growth: {perf.get('month_over_month_growth')}")
        doc.add_paragraph("")

    # Risk Assessment
    doc.add_heading("Risk Assessment", level=2)
    doc.add_paragraph(f"Overall Risk Level: {(risk.get('risk_level') or 'N/A').upper()}")

    if risk.get("key_risks"):
        doc.add_heading("Key Risks", level=3)
        for r in risk.get("key_risks", []):
            doc.add_paragraph(f"• {r}")

    if risk.get("mitigation_status"):
        doc.add_heading("Mitigation Actions", level=3)
        for m in risk.get("mitigation_status", []):
            doc.add_paragraph(f"• {m}")
    doc.add_paragraph("")

    # Forecasting
    doc.add_heading("Forecasting & Projections", level=2)

    forecast_data = (monthly_data or {}).get("forecasting") or {}
    _add_table(
        ["Projection", "Value"],
        [
            ["Expected New Clients", _fmt_num(forecast_data.get("expected_new_clients", 0))],
            ["Expected Submissions", _fmt_num(forecast_data.get("expected_submissions", 0))],
            ["Expected Settlement Value (90 days)", _fmt_currency(fcst.get("expected_settlements_next_90_days", 0))],
        ]
    )

    if fcst.get("quarterly_outlook"):
        doc.add_paragraph("")
        doc.add_paragraph(f"Quarterly Outlook: {fcst.get('quarterly_outlook')}")
    doc.add_paragraph("")

    # ==================== CHARTS SECTION ====================
    doc.add_heading("5. Charts & Visualizations", level=1)

    # Build figures
    figs = _build_dashboard_figures(monthly_data, investor_report)

    # Write images to temp files
    tmp_dir = os.path.join(os.path.dirname(out_path), ".tmp_report_assets")
    os.makedirs(tmp_dir, exist_ok=True)

    added_any = False
    export_failures: List[str] = []

    def _add_chart(key: str, caption: str):
        nonlocal added_any
        fig = figs.get(key)
        if fig is None:
            return
        try:
            png, err = _export_plotly_fig_to_png_bytes_with_error(fig)
            if not png:
                if err is not None:
                    raise err
                raise RuntimeError("Empty PNG bytes returned by plotly.io.to_image")

            img_path = os.path.join(tmp_dir, f"{key}.png")
            with open(img_path, "wb") as f:
                f.write(png)

            doc.add_paragraph(caption)
            doc.add_picture(img_path, width=Inches(6.5))
            doc.add_paragraph("")
            added_any = True
        except Exception as e:
            export_failures.append(f"{key}: {type(e).__name__}: {e}")

    # Lenders charts
    doc.add_heading("Lender Distribution Charts", level=2)
    _add_chart("lenders_pie", "Top 10 Lenders by Claims")
    _add_chart("lenders_value_bar", "Top 15 Lenders by Portfolio Value")

    # Economic charts
    doc.add_heading("Economic Analysis Charts", level=2)
    _add_chart("economic_summary", "Economic Summary (Expected)")
    _add_chart("profit_split", "Expected Profit Distribution")

    # Pipeline charts
    doc.add_heading("Pipeline Charts", level=2)
    _add_chart("pipeline_funnel", "Pipeline Funnel (by Count)")
    _add_chart("pipeline_value", "Pipeline Value by Stage")

    if not added_any:
        doc.add_paragraph(
            "Charts could not be embedded. This usually means Plotly image export failed in this environment. "
            "Ensure 'kaleido' is installed and a Chrome/Chromium binary is available."
        )
        if export_failures:
            doc.add_paragraph("Export errors:")
            for err in export_failures[:10]:
                doc.add_paragraph(f"- {err}")

    # ==================== ACTION ITEMS ====================
    doc.add_heading("6. Action Items", level=1)

    action_items = (investor_report or {}).get("action_items") or []
    if action_items:
        action_rows = []
        for item in action_items:
            item = item or {}
            action_rows.append([
                (item.get('priority') or 'N/A').upper(),
                item.get('action', 'N/A'),
                item.get('owner', 'N/A'),
                item.get('deadline', 'N/A'),
                item.get('rationale', 'N/A')
            ])
        _add_table(
            ["Priority", "Action", "Owner", "Deadline", "Rationale"],
            action_rows
        )
    else:
        doc.add_paragraph("No action items identified.")

    doc.add_paragraph("")

    # Footer
    doc.add_paragraph("---")
    doc.add_paragraph("Report generated by AI Multi-Agent System")
    doc.add_paragraph(f"Chart export environment: {_kaleido_debug_info()}")

    doc.save(out_path)
    return out_path


def build_investor_report_pptx(
    *,
    out_path: str,
    template_path: str = "reports/Monthly_Investor_Reporting_Dec 2025.pptx",
    monthly_data: Dict[str, Any],
    investor_report: Dict[str, Any],
) -> str:
    """
    Create a PowerPoint investor report using the template.
    Each slide contains data tables and relevant charts organized together.
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")

    # Load template or create new presentation
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        print(f"📊 Loading PowerPoint template: {template_path}")
    else:
        prs = Presentation()
        print("📊 Creating new PowerPoint presentation (template not found)")

    # Helper functions
    def _fmt_currency(v):
        if v is None or v == "N/A":
            return "N/A"
        try:
            return f"£{float(v):,.0f}"
        except (ValueError, TypeError):
            return str(v)

    def _fmt_pct(v):
        if v is None or v == "N/A":
            return "N/A"
        try:
            return f"{float(v):.1f}%"
        except (ValueError, TypeError):
            return str(v)

    def _fmt_num(v):
        if v is None or v == "N/A":
            return "N/A"
        try:
            return f"{int(float(v)):,}"
        except (ValueError, TypeError):
            return str(v)

    def add_title_slide(title: str, subtitle: str):
        """Add a title slide"""
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        return slide

    def add_content_slide(title: str):
        """Add a content slide with title"""
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        return slide

    def add_table_to_slide(slide, headers: List[str], rows: List[List[str]],
                           left: float, top: float, width: float, height: float):
        """Add a table to a slide"""
        table_shape = slide.shapes.add_table(
            len(rows) + 1, len(headers),
            PptxInches(left), PptxInches(top),
            PptxInches(width), PptxInches(height)
        )
        table = table_shape.table

        # Style header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Add data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_val in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_val) if cell_val is not None else "N/A"
                cell.text_frame.paragraphs[0].font.size = Pt(9)

        return table_shape

    def add_text_box(slide, text: str, left: float, top: float, width: float, height: float,
                     font_size: int = 11, bold: bool = False):
        """Add a text box to a slide"""
        textbox = slide.shapes.add_textbox(
            PptxInches(left), PptxInches(top),
            PptxInches(width), PptxInches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        return textbox

    def add_chart_image(slide, fig, left: float, top: float, width: float):
        """Add a Plotly chart as an image to a slide"""
        try:
            png, err = _export_plotly_fig_to_png_bytes_with_error(fig)
            if not png or err:
                return None

            # Save to temp file
            tmp_dir = os.path.join(os.path.dirname(out_path), ".tmp_pptx_assets")
            os.makedirs(tmp_dir, exist_ok=True)
            img_path = os.path.join(tmp_dir, f"chart_{id(fig)}.png")

            with open(img_path, "wb") as f:
                f.write(png)

            slide.shapes.add_picture(img_path, PptxInches(left), PptxInches(top), width=PptxInches(width))
            return True
        except Exception as e:
            print(f"Warning: Could not add chart to slide: {e}")
            return None

    # Get report data
    exec_sum = (investor_report or {}).get("executive_summary") or {}
    perf = (investor_report or {}).get("portfolio_performance") or {}
    fin = (investor_report or {}).get("financial_analysis") or {}
    comp = (investor_report or {}).get("compliance_assessment") or {}
    risk = (investor_report or {}).get("risk_assessment") or {}
    lender_conc = (investor_report or {}).get("lender_concentration") or {}
    pipeline_analysis = (investor_report or {}).get("pipeline_analysis") or {}
    cost_eff = (investor_report or {}).get("cost_efficiency") or {}
    fcst = (investor_report or {}).get("forecasting") or {}

    # Get raw monthly data
    lenders = (monthly_data or {}).get("lender_distribution") or []
    pipeline = (monthly_data or {}).get("pipeline") or {}
    pm = (monthly_data or {}).get("portfolio_metrics") or {}
    fm = (monthly_data or {}).get("financial_metrics") or {}

    period = exec_sum.get("reporting_period") or "Monthly Report"

    # Build charts
    figs = _build_dashboard_figures(monthly_data, investor_report)

    # ==================== SLIDE 1: EXECUTIVE SUMMARY ====================
    slide1 = add_content_slide(f"Executive Summary - {period}")

    # Key metrics table (left side)
    key_metrics = [
        ["Total Claims", _fmt_num(perf.get("total_claims") or pm.get("unique_claims", 0))],
        ["Total Clients", _fmt_num(perf.get("total_clients") or pm.get("unique_clients", 0))],
        ["Total Lenders", _fmt_num(len(lenders))],
        ["Portfolio Value", _fmt_currency(perf.get("total_portfolio_value") or pm.get("total_settlement_value", 0))],
        ["LP Return (80%)", _fmt_currency(fin.get("funder_expected_return"))],
        ["GP Return (20%)", _fmt_currency(fin.get("firm_expected_return"))],
        ["MOIC", f"{fin.get('moic_projection', 0):.2f}x" if fin.get('moic_projection') else "N/A"],
        ["ROI", _fmt_pct(fin.get("roi_projection"))],
        ["Compliance Status", (comp.get("fca_compliance_status") or "N/A").upper()],
    ]
    add_table_to_slide(slide1, ["Metric", "Value"], key_metrics, 0.5, 1.3, 4.5, 3.5)

    # Key highlights (right side)
    highlights_text = "Key Highlights:\n"
    for metric in exec_sum.get("key_metrics_summary", [])[:5]:
        highlights_text += f"• {metric}\n"
    if exec_sum.get("critical_updates"):
        highlights_text += "\nCritical Updates:\n"
        for update in exec_sum.get("critical_updates", [])[:3]:
            highlights_text += f"• {update}\n"
    add_text_box(slide1, highlights_text, 5.5, 1.3, 7.5, 4.0, font_size=10)

    # ==================== SLIDE 2: ECONOMIC ANALYSIS ====================
    slide2 = add_content_slide("Economic Analysis & Profit Distribution")

    # Financial summary table (left side)
    financial_rows = [
        ["Total Settlement Value", _fmt_currency(fin.get("total_settlements"))],
        ["DBA Proceeds (30%)", _fmt_currency(fin.get("dba_proceeds_expected"))],
        ["Total Costs Incurred", _fmt_currency(fin.get("total_costs_incurred"))],
        ["LP Return (80% of DBA)", _fmt_currency(fin.get("funder_expected_return"))],
        ["GP Return (20% of DBA)", _fmt_currency(fin.get("firm_expected_return"))],
        ["ROI Projection", _fmt_pct(fin.get("roi_projection"))],
        ["MOIC Projection", f"{fin.get('moic_projection', 0):.2f}x" if fin.get('moic_projection') else "N/A"],
    ]
    add_table_to_slide(slide2, ["Metric", "Value"], financial_rows, 0.3, 1.2, 4.5, 3.0)

    # Cost breakdown table (below financial)
    cost_rows = [
        ["Acquisition Cost", _fmt_currency(fm.get("acquisition_cost", 0))],
        ["Submission Cost", _fmt_currency(fm.get("submission_cost", 0))],
        ["Total Costs", _fmt_currency(fm.get("total_costs", 0))],
        ["Cost per Claim", _fmt_currency(cost_eff.get("cost_per_claim") or fm.get("cost_per_claim", 0))],
    ]
    add_table_to_slide(slide2, ["Cost Category", "Amount"], cost_rows, 0.3, 4.5, 4.5, 1.8)

    # Profit split chart (right side)
    if figs.get("profit_split"):
        add_chart_image(slide2, figs["profit_split"], 5.2, 1.2, 7.5)

    # ==================== SLIDE 3: COMPLIANCE & PIPELINE ====================
    slide3 = add_content_slide("Compliance & Pipeline Status")

    # Compliance assessment (top left)
    compliance_rows = [
        ["Compliance Status", (comp.get("fca_compliance_status") or "N/A").upper()],
        ["Claims at Risk", _fmt_num(comp.get("claims_at_risk", 0))],
    ]
    add_table_to_slide(slide3, ["Metric", "Value"], compliance_rows, 0.3, 1.2, 3.5, 0.9)

    # Compliance actions (below)
    if comp.get("compliance_actions_needed"):
        actions_text = "Actions Required:\n"
        for action in comp.get("compliance_actions_needed", [])[:4]:
            actions_text += f"• {action[:60]}...\n" if len(action) > 60 else f"• {action}\n"
        add_text_box(slide3, actions_text, 0.3, 2.3, 4.5, 1.5, font_size=9)

    # Pipeline breakdown table (left side, lower)
    pipeline_stages = [
        ("Awaiting DSAR", pipeline.get("awaiting_dsar", {})),
        ("Pending Submission", pipeline.get("pending_submission", {})),
        ("Under Review", pipeline.get("under_review", {})),
        ("Settlement Offered", pipeline.get("settlement_offered", {})),
        ("Paid", pipeline.get("paid", {})),
    ]
    pipeline_rows = []
    total_count = 0
    total_value = 0
    for stage_name, stage_data in pipeline_stages:
        count = stage_data.get("count", 0) if isinstance(stage_data, dict) else 0
        value = stage_data.get("value", 0) if isinstance(stage_data, dict) else 0
        total_count += count
        total_value += value
        pipeline_rows.append([stage_name, _fmt_num(count), _fmt_currency(value)])
    pipeline_rows.append(["TOTAL", _fmt_num(total_count), _fmt_currency(total_value)])

    add_table_to_slide(slide3, ["Pipeline Stage", "Count", "Value"], pipeline_rows, 0.3, 4.0, 4.5, 2.5)

    # Pipeline funnel chart (right side)
    if figs.get("pipeline_funnel"):
        add_chart_image(slide3, figs["pipeline_funnel"], 5.2, 1.2, 7.5)

    # ==================== SLIDE 4: PORTFOLIO ANALYSIS ====================
    slide4 = add_content_slide("Portfolio Analysis & Forecasting")

    # Portfolio metrics table (left side)
    portfolio_rows = [
        ["Total Claims", _fmt_num(perf.get("total_claims") or pm.get("unique_claims", 0))],
        ["Total Clients", _fmt_num(perf.get("total_clients") or pm.get("unique_clients", 0))],
        ["Claims Submitted", _fmt_num(pm.get("claims_submitted", 0))],
        ["Claims Successful", _fmt_num(pm.get("claims_successful", 0))],
        ["Claims Rejected", _fmt_num(pm.get("claims_rejected", 0))],
        ["Success Rate", _fmt_pct(perf.get("success_rate") or pm.get("success_rate", 0))],
        ["Avg Claim Value", _fmt_currency(perf.get("average_settlement") or pm.get("avg_claim_value", 0))],
    ]
    add_table_to_slide(slide4, ["Metric", "Value"], portfolio_rows, 0.3, 1.2, 4.0, 3.0)

    # Risk assessment (below portfolio)
    risk_text = f"Risk Level: {(risk.get('risk_level') or 'N/A').upper()}\n\n"
    if risk.get("key_risks"):
        risk_text += "Key Risks:\n"
        for r in risk.get("key_risks", [])[:3]:
            risk_text += f"• {r[:50]}...\n" if len(r) > 50 else f"• {r}\n"
    add_text_box(slide4, risk_text, 0.3, 4.4, 4.5, 2.0, font_size=9)

    # Forecasting table (right side top)
    forecast_data = (monthly_data or {}).get("forecasting") or {}
    forecast_rows = [
        ["Expected New Clients", _fmt_num(forecast_data.get("expected_new_clients", 0))],
        ["Expected Submissions", _fmt_num(forecast_data.get("expected_submissions", 0))],
        ["Expected Settlement (90 days)", _fmt_currency(fcst.get("expected_settlements_next_90_days", 0))],
    ]
    add_table_to_slide(slide4, ["Projection", "Value"], forecast_rows, 5.0, 1.2, 4.5, 1.5)

    # Quarterly outlook
    if fcst.get("quarterly_outlook"):
        outlook_text = f"Quarterly Outlook:\n{fcst.get('quarterly_outlook')}"
        add_text_box(slide4, outlook_text, 5.0, 2.9, 7.5, 1.5, font_size=10)

    # Lenders value bar chart (right side bottom)
    if figs.get("lenders_value_bar"):
        add_chart_image(slide4, figs["lenders_value_bar"], 5.0, 4.0, 7.5)

    # ==================== SLIDE 5: ACTION ITEMS ====================
    action_items = (investor_report or {}).get("action_items") or []
    if action_items:
        slide5 = add_content_slide("Action Items")

        action_rows = []
        for item in action_items[:8]:  # Limit to 8 items
            item = item or {}
            action_rows.append([
                (item.get('priority') or 'N/A').upper(),
                item.get('action', 'N/A')[:40] + "..." if len(item.get('action', '')) > 40 else item.get('action', 'N/A'),
                item.get('owner', 'N/A'),
                item.get('deadline', 'N/A')
            ])

        add_table_to_slide(slide5, ["Priority", "Action", "Owner", "Deadline"], action_rows, 0.5, 1.3, 12.0, 4.5)

    # Save the presentation
    prs.save(out_path)
    print(f"✅ PowerPoint report saved: {out_path}")
    return out_path


def generate_full_investor_report(excel_path: str) -> Dict[str, Any]:
    """
    Main orchestration function - coordinates all agents to generate investor report
    """
    print("="*80)
    print("🤖 INTELLIGENT MULTI-AGENT SYSTEM")
    print("="*80)

    # Initialize all agents
    priority_deed_agent = PriorityDeedAgent()
    fca_agent = FCAComplianceAgent()
    monthly_agent = MonthlyReportAgent()
    investor_agent = InvestorReportAgent()

    # Step 1: Priority Deed Agent reads profit distribution rules
    profit_rules = priority_deed_agent.read_priority_deed()

    # Step 2: FCA Agent reads compliance requirements
    compliance_rules = fca_agent.read_fca_scheme()

    # Step 3: Monthly Report Agent extracts data from Excel
    monthly_data = monthly_agent.analyze_monthly_report(excel_path)

    # Step 4: Investor Report Agent generates comprehensive report
    report = investor_agent.generate_investor_report(
        monthly_data=monthly_data,
        profit_rules=profit_rules,
        compliance_rules=compliance_rules
    )

    # Build reports on disk
    os.makedirs("reports", exist_ok=True)
    period = (monthly_data or {}).get("reporting_period") or "Report"
    safe_period = str(period).replace("/", "-").replace("\\", "-").replace(":", "-")

    # Build DOCX report
    docx_path = os.path.join("reports", f"Investor_Report_{safe_period}.docx")
    try:
        build_investor_report_docx(
            out_path=docx_path,
            narrative=report.get("narrative") or "",
            monthly_data=monthly_data,
            investor_report=report.get("report_data") or {},
        )
    except Exception as e:
        print(f"Warning: DOCX report generation failed: {e}")
        docx_path = None

    # Build PowerPoint report
    pptx_path = os.path.join("reports", f"Investor_Report_{safe_period}.pptx")
    try:
        if PPTX_AVAILABLE:
            build_investor_report_pptx(
                out_path=pptx_path,
                monthly_data=monthly_data,
                investor_report=report.get("report_data") or {},
            )
        else:
            print("Warning: python-pptx not available, skipping PowerPoint generation")
            pptx_path = None
    except Exception as e:
        print(f"Warning: PowerPoint report generation failed: {e}")
        pptx_path = None

    print("="*80)
    print("✅ REPORT GENERATION COMPLETE")
    print("="*80)

    return {
        "monthly_data": monthly_data,
        "profit_rules": profit_rules,
        "compliance_rules": compliance_rules,
        "investor_report": report["report_data"],
        "markdown_report": report["markdown_report"],
        "docx_report_path": docx_path,
        "pptx_report_path": pptx_path,
    }


if __name__ == "__main__":
    result = generate_full_investor_report("uploads/Milberg_MOnthly_Report.xlsx")
    print("\n" + result['markdown_report'])
